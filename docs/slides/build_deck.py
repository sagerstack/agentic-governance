"""Generate the 3-slide Agentic Governance deck from the research + gap artifacts."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ----
INK      = RGBColor(0x0F, 0x17, 0x2A)  # deep navy
ACCENT   = RGBColor(0x2D, 0x6C, 0xDF)  # blue
ACCENT2  = RGBColor(0x12, 0x9E, 0x7F)  # teal/green
DANGER   = RGBColor(0xC0, 0x3A, 0x2B)  # red
AMBER    = RGBColor(0xB8, 0x7A, 0x00)  # amber
MUTED    = RGBColor(0x5B, 0x66, 0x77)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xF2, 0xF5, 0xFA)
HEADROW  = RGBColor(0x1B, 0x2A, 0x45)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def band(slide, title, subtitle, accent=ACCENT):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    _set_fill(bar, INK)
    strip = slide.shapes.add_shape(1, 0, Inches(1.15), SW, Inches(0.09))
    _set_fill(strip, accent)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.16), SW - Inches(1.1), Inches(0.95)).text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tb.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xC6, 0xD3, 0xEA)


def textbox(slide, l, t, w, h):
    tf = slide.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=13, bold=False, color=INK, bullet=False, space=4, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space)
    r = p.add_run(); r.text = ("•  " if bullet else "") + text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return p


def card(slide, l, t, w, h, fill=LIGHT):
    c = slide.shapes.add_shape(1, l, t, w, h)
    _set_fill(c, fill)
    c.shadow.inherit = False
    return c


def table(slide, rows, cols, l, t, w, h, col_w=None):
    gt = slide.shapes.add_table(rows, cols, l, t, w, h).table
    if col_w:
        for i, cw in enumerate(col_w):
            gt.columns[i].width = cw
    return gt


def style_header(gt, headers, size=12):
    for j, htext in enumerate(headers):
        cell = gt.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = HEADROW
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = htext
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = WHITE


def cell_text(gt, i, j, text, size=11, bold=False, color=INK, fill=None, align=PP_ALIGN.LEFT):
    cell = gt.cell(i, j)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color


# ============================================================= SLIDE 1
s1 = prs.slides.add_slide(BLANK)
band(s1, "AI Governance Standards",
     "Research summary — which frameworks we reviewed and what they mandate vs. recommend", ACCENT)

# two standard cards
cw = Inches(6.0); ch = Inches(3.1); top = Inches(1.5)
c1 = card(s1, Inches(0.55), top, cw, ch, LIGHT)
t1 = textbox(s1, Inches(0.8), top + Inches(0.18), cw - Inches(0.5), ch - Inches(0.3))
para(t1, "1 · Singapore IMDA — Model AI Governance Framework", 15, True, ACCENT, first=True)
para(t1, "Voluntary “model” frameworks: 2020 MGF · Generative AI (2024, 9 dimensions) · Agentic AI (2026)", 11.5, color=MUTED)
para(t1, "Agent identity & authorization; least-privilege > prompt-layer", 12, bullet=True)
para(t1, "MCP as a governance layer; input/output filters; RAG grounding", 12, bullet=True)
para(t1, "Risk-based human oversight; monitoring, failsafes, incident reporting", 12, bullet=True)

c2 = card(s1, Inches(6.78), top, cw, ch, LIGHT)
t2 = textbox(s1, Inches(7.03), top + Inches(0.18), cw - Inches(0.5), ch - Inches(0.3))
para(t2, "2 · MAS — Safeguards for Agentic Finance at Runtime (SAFR)", 15, True, ACCENT2, first=True)
para(t2, "White paper v1.0, 3 Jul 2026  ·  + MAS FEAT / Veritas (analogues)", 11.5, color=MUTED)
para(t2, "Purpose-built for runtime, pre-execution control of agent actions", 12, bullet=True)
para(t2, "4-component checkpoint: Identity → Controls → Disposition → Audit", 12, bullet=True)
para(t2, "Governance Envelope → Deny / Escalate / Auto-Execute / Observe", 12, bullet=True)

# mandatory vs recommended strip
mb = card(s1, Inches(0.55), Inches(4.78), Inches(12.23), Inches(2.15), RGBColor(0xFF, 0xF4, 0xE0))
mt = textbox(s1, Inches(0.8), Inches(4.9), Inches(11.9), Inches(1.95))
para(mt, "Mandatory vs. Recommended — key finding", 14, True, AMBER, first=True)
para(mt, "No legal mandate. Both are advisory/voluntary — SAFR states it “does not constitute regulatory guidance”; IMDA is a voluntary model framework; MAS AI Risk Mgmt Guidelines are still in consultation.", 12, color=INK)
para(mt, "Classified on two axes → Binding force = Recommended (advisory) for ALL controls · Source strength = core / should / may. “Mandatory” in our catalogue = a framework’s own core/“must” language, not law.", 12, color=INK)
para(mt, "Genuine legal mandates would arise only from applicable law + a finalised MAS AIRG (out of scope).", 11.5, color=MUTED)

# ============================================================= SLIDE 2
s2 = prs.slides.add_slide(BLANK)
band(s2, "Agentic Defence",
     "Four control groups form a defence-in-depth control plane — what each defends, what we build, and the risk it retires", ACCENT2)

intro = textbox(s2, Inches(0.55), Inches(1.4), Inches(12.2), Inches(0.42))
para(intro, "Groups A–D are the four layers of the governance control plane. Each attaches to a specific seam in the agentic app.",
     12.5, color=INK, first=True)

gt = table(s2, 5, 4, Inches(0.5), Inches(1.9), Inches(12.33), Inches(4.05),
           col_w=[Inches(3.5), Inches(3.85), Inches(2.35), Inches(2.63)])
style_header(gt, ["Group — what it defends", "What we build", "Aligns to", "Protects against"], size=11.5)

groups = [
 ("A", "Action-time authorization", "Gate every tool action before it executes (the SAFR checkpoint)", DANGER,
  "Envelope from trusted state · agent identity · mandate · deny / escalate / auto disposition · exposure & rate limits — at mcpCallTool()",
  "MAS SAFR 4-component checkpoint · IMDA least-privilege",
  "Fraudulent payout · hijacked / unauthorized actions · action abuse"),
 ("B", "Model input/output guardrails", "Inspect what enters and leaves the LLM / VLM", ACCENT,
  "Prompt-injection detection on chat + receipt image · PII redaction · evidence-grounded output — at model hooks",
  "IMDA GenAI input filters · grounded output / RAG",
  "Agent hijack via malicious receipt / injection · PII leakage"),
 ("C", "Human oversight & failsafes", "Keep a human in control of high-risk actions", AMBER,
  "Risk-calibrated escalation · timeout → default-deny contract · fail-closed floor · employee recourse — at humanEscalation",
  "SAFR human-reviewer escalation · IMDA risk-based oversight",
  "Unreviewed high-risk / irreversible actions · rubber-stamp review"),
 ("D", "Audit, monitoring & incident", "Tamper-evident record + detect and respond", ACCENT2,
  "Immutable, PII-safe audit log · real-time monitoring → intervention · incident workflow — at the audit log",
  "SAFR Audit Log · IMDA tamper-evident log & monitoring",
  "Audit tampering · undetected attacks · forensic gaps"),
]
for i, (letter, name, defn, col, build, aligns, protects) in enumerate(groups, start=1):
    # column 0: rich group cell (letter+name bold colored, definition muted)
    c0 = gt.cell(i, 0)
    c0.fill.solid(); c0.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
    c0.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf0 = c0.text_frame; tf0.word_wrap = True
    p0 = tf0.paragraphs[0]; r0 = p0.add_run(); r0.text = f"{letter} · {name}"
    r0.font.size = Pt(12); r0.font.bold = True; r0.font.color.rgb = col
    p0b = tf0.add_paragraph(); r0b = p0b.add_run(); r0b.text = defn
    r0b.font.size = Pt(9.5); r0b.font.color.rgb = MUTED
    cell_text(gt, i, 1, build, 9.5)
    cell_text(gt, i, 2, aligns, 9.5)
    cell_text(gt, i, 3, protects, 9.5)

foot = card(s2, Inches(0.5), Inches(6.15), Inches(12.33), Inches(1.02), INK)
ft = textbox(s2, Inches(0.75), Inches(6.24), Inches(11.9), Inches(0.85))
para(ft, "What we’re building now:  Group A first — the highest attack-defense value (83% of it is absent today).  An app-agnostic, in-process control plane at the tool boundary, using deterministic checks against forge-proof trusted state (no LLM in the loop).",
     11.5, True, WHITE, first=True)
para(ft, "Group A POC = 6 thin slices: envelope → deny-unknown-tool → identity + mandate → integrity → exposure / rate / evidence → schema + escalate.  Groups B, C, D follow incrementally.",
     10.5, color=RGBColor(0xC6, 0xD3, 0xEA))

# ============================================================= SLIDE 3
s3 = prs.slides.add_slide(BLANK)
band(s3, "Gap Analysis — Current Expense AI App",
     "Evidence-based code audit — grouped A–D, scored by enforcement strength", DANGER)

hl = card(s3, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.6), RGBColor(0xFD, 0xEC, 0xEA))
ht = textbox(s3, Inches(0.75), Inches(1.44), Inches(11.9), Inches(0.52))
para(ht, "Strong seams, near-zero enforcement — 0 of 28 controls are actually enforced today. The interception points exist, but none governs.",
     13, True, DANGER, first=True)

# scoring legend
lg = card(s3, Inches(0.5), Inches(2.12), Inches(12.33), Inches(0.92), LIGHT)
lt = textbox(s3, Inches(0.75), Inches(2.19), Inches(11.9), Inches(0.82))
p = lt.paragraphs[0]; p.space_after = Pt(3)
for label, defn, col in [("Present", " = built AND deterministically enforced.   ", RGBColor(0x2B,0x8A,0x3E)),
                          ("Partial", " = a mechanism exists but is prompt-based, incomplete or non-enforcing.   ", AMBER),
                          ("Absent", " = none found.", DANGER)]:
    r = p.add_run(); r.text = label; r.font.bold = True; r.font.size = Pt(11); r.font.color.rgb = col
    r2 = p.add_run(); r2.text = defn; r2.font.size = Pt(11); r2.font.color.rgb = INK
pc = lt.add_paragraph()
rr = pc.add_run(); rr.text = "Scored today:  0 Present  ·  11 Partial  ·  17 Absent   (of 28 controls)   →   61% missing (Absent),  0% enforced"
rr.font.bold = True; rr.font.size = Pt(11.5); rr.font.color.rgb = DANGER

# per-group detail table
gt3 = table(s3, 5, 4, Inches(0.5), Inches(3.18), Inches(12.33), Inches(3.0),
            col_w=[Inches(3.15), Inches(2.55), Inches(3.35), Inches(3.28)])
style_header(gt3, ["Group — what it defends", "Missing (Absent)", "What exists today (Partial)", "Biggest gap (Absent)"], size=11)
grows = [
 ("A", "Action-time authorization", "gate tool actions", DANGER,
  "83% missing", "10 of 12 absent · 2 partial · 0 enforced",
  "Typed MCP schemas; a SELECT-only query guard; a confidence threshold that is dead config",
  "No authz at mcpCallTool — no envelope, identity, mandate, disposition, or value/rate limits"),
 ("B", "Model input/output guardrails", "screen LLM/VLM I/O", ACCENT,
  "33% missing", "2 of 6 absent · 4 partial · 0 enforced",
  "Grounding, PII handling & graceful-failure exist — but prompt-based / model-asserted / log-only",
  "No prompt-injection detection; the receipt image is trusted as input"),
 ("C", "Human oversight & failsafes", "human-in-control", AMBER,
  "20% missing", "1 of 5 absent · 4 partial · 0 enforced",
  "Real reviewer authority; escalation persists a status; some conservative error-escalation",
  "No timeout / deadline / default-deny; no oversight-effectiveness metrics"),
 ("D", "Audit, monitoring & incident", "record & detect", ACCENT2,
  "80% missing", "4 of 5 absent · 1 partial · 0 enforced",
  "End-to-end trace across logs + LangGraph checkpointer (correlated by claim)",
  "Audit log is mutable & CASCADE-deletable; no real-time monitoring→intervention or incident flow"),
]
for i, (letter, name, defn, col, enf, split, exists, gap) in enumerate(grows, start=1):
    c0 = gt3.cell(i, 0)
    c0.fill.solid(); c0.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
    c0.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf0 = c0.text_frame; tf0.word_wrap = True
    p0 = tf0.paragraphs[0]; r0 = p0.add_run(); r0.text = f"{letter} · {name}"
    r0.font.size = Pt(11.5); r0.font.bold = True; r0.font.color.rgb = col
    p0b = tf0.add_paragraph(); r0b = p0b.add_run(); r0b.text = defn
    r0b.font.size = Pt(9.5); r0b.font.color.rgb = MUTED
    # missing % cell (big % + breakdown)
    c1 = gt3.cell(i, 1)
    c1.fill.solid(); c1.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
    c1.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf1 = c1.text_frame; tf1.word_wrap = True
    q0 = tf1.paragraphs[0]; e0 = q0.add_run(); e0.text = enf
    e0.font.size = Pt(15); e0.font.bold = True; e0.font.color.rgb = col
    q1 = tf1.add_paragraph(); e1 = q1.add_run(); e1.text = split
    e1.font.size = Pt(9); e1.font.color.rgb = MUTED
    cell_text(gt3, i, 2, exists, 9.5)
    cell_text(gt3, i, 3, gap, 9.5, color=RGBColor(0x8A,0x2A,0x1E))

concl = card(s3, Inches(0.5), Inches(6.35), Inches(12.33), Inches(0.72), INK)
ct = textbox(s3, Inches(0.75), Inches(6.42), Inches(11.9), Inches(0.6))
para(ct, "Overall: 61% of controls are missing entirely (17 of 28 Absent) and 0% are enforced.  Group A is the worst — 10 of 12 absent — so the POC targets it first.",
     12, True, WHITE, first=True)

# ============================================================= SLIDE 4
s4 = prs.slides.add_slide(BLANK)
band(s4, "Building agentic-governance — an app-agnostic control plane",
     "A pure decision core + injected adapters; the app integrates at one seam", ACCENT)


def lbox(slide, l, t, w, h, title, subs, fill, tcolor=WHITE, scolor=None):
    b = slide.shapes.add_shape(1, l, t, w, h); _set_fill(b, fill)
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.07)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.bold = True; r.font.size = Pt(12.5); r.font.color.rgb = tcolor
    for s in subs:
        pp = tf.add_paragraph(); pp.space_before = Pt(1); rr = pp.add_run(); rr.text = s
        rr.font.size = Pt(9.5); rr.font.color.rgb = scolor or tcolor
    return b


def arrow(slide, cx, t):
    a = slide.shapes.add_textbox(cx - Inches(0.25), t, Inches(0.5), Inches(0.24)).text_frame
    a.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = a.paragraphs[0].add_run(); r.text = "▼"; r.font.size = Pt(14); r.font.color.rgb = MUTED

LX = Inches(0.5); LW = Inches(6.7); cx = Inches(0.5) + Inches(3.35)
# layer 3 - host app
lbox(s4, LX, Inches(1.5), LW, Inches(0.92),
     "Host agentic app  (e.g. Expense AI)  — unchanged logic",
     ["Sets trusted state via contextvars · one minimal DI edit · calls the tool boundary"],
     RGBColor(0x33, 0x43, 0x63), WHITE, RGBColor(0xCF, 0xDA, 0xEE))
arrow(s4, cx, Inches(2.44))
# layer 2 - integration adapter
lbox(s4, LX, Inches(2.7), LW, Inches(0.98),
     "Integration adapter  —  app glue (one small module per app)",
     ["governedMcpCallTool wrapper (identical signature)  ·  CallContext assembler",
      "Injected, NOT imported — receives providers + real tool via install()"],
     ACCENT, WHITE, RGBColor(0xE6, 0xEE, 0xFB))
arrow(s4, cx, Inches(3.7))
# layer 1 - core container
core = card(s4, LX, Inches(3.96), LW, Inches(2.92), INK)
ctf = textbox(s4, LX + Inches(0.16), Inches(4.02), LW - Inches(0.32), Inches(0.5))
para(ctf, "Governance CORE  —  app-agnostic (imports nothing app-specific)", 12.5, True, WHITE, first=True)
chips = [
  ("Decision engine  →  Deny · Escalate · Auto-Execute · Observe", "+ Governance Envelope built from trusted state"),
  ("Ports (interfaces)", "identity · mandate · policy-decision-point · counters · evidence · audit"),
  ("Adapters (pluggable)", "pure-Python PDP · JSONL audit sink · in-memory registry & counters"),
]
cy = Inches(4.5)
for head, sub in chips:
    ch = card(s4, LX + Inches(0.2), cy, LW - Inches(0.4), Inches(0.68), RGBColor(0x22, 0x31, 0x4E))
    ctf2 = textbox(s4, LX + Inches(0.34), cy + Inches(0.04), LW - Inches(0.7), Inches(0.62))
    para(ctf2, head, 10.5, True, WHITE, first=True, space=0)
    para(ctf2, sub, 9.5, False, RGBColor(0xB8, 0xC6, 0xE0), space=0)
    cy = cy + Inches(0.76)

# right column
RX = Inches(7.42); RW = Inches(5.4)
rt = textbox(s4, RX, Inches(1.5), RW, Inches(2.7))
para(rt, "App-agnostic by construction", 15, True, ACCENT, first=True)
para(rt, "core/ and ports/ import nothing app-specific (stdlib only) — verified in review", 11.5, bullet=True, color=INK)
para(rt, "The app injects its providers + real tool via install() — inject, not import; the core has zero dependency on the app", 11.5, bullet=True, color=INK)
para(rt, "Policies (allowlists, mandates, thresholds) live in the governance repo, not the app", 11.5, bullet=True, color=INK)
para(rt, "One integration seam: the tool boundary (mcpCallTool); model hooks added later for Group B", 11.5, bullet=True, color=INK)

ob = card(s4, RX, Inches(4.35), RW, Inches(1.55), RGBColor(0xE7, 0xF4, 0xEF))
ot = textbox(s4, RX + Inches(0.22), Inches(4.45), RW - Inches(0.44), Inches(1.4))
para(ot, "Onboard any agentic app in 2 steps", 13, True, ACCENT2, first=True)
para(ot, "1 · Write a thin adapter mapping its ambient state → CallContext", 11, color=INK)
para(ot, "2 · Add one composition-root DI edit to call the wrapper", 11, color=INK)
para(ot, "→ Core + policy engine unchanged.", 11, bold=True, color=ACCENT2)

fb = card(s4, RX, Inches(6.05), RW, Inches(0.82), INK)
ft4 = textbox(s4, RX + Inches(0.22), Inches(6.12), RW - Inches(0.44), Inches(0.7))
para(ft4, "Deterministic decisions · no LLM in the loop · fail-closed when governance is unavailable", 11, True, WHITE, first=True)

# ============================================================= SPEAKER NOTES
def set_notes(slide, paras):
    tf = slide.notes_slide.notes_text_frame
    tf.text = paras[0]
    for p in paras[1:]:
        para = tf.add_paragraph(); para.text = p

set_notes(s1, [
  "PURPOSE OF THIS SLIDE: establish which external AI-governance standards we reviewed to design our controls, and clarify their legal status (advisory, not law). Everything we build traces back to these two families.",
  "STANDARD 1 — Singapore IMDA Model AI Governance Framework. This is a family of voluntary 'model' frameworks published by IMDA / AI Verify Foundation: the 2020 Model AI Governance Framework (2nd ed.), the 2024 Generative AI edition (nine governance dimensions), and the 2026 Agentic AI edition. For our purposes the useful runtime content is: verifiable agent identity and authorization; the principle that least-privilege and controls must be enforced deterministically at the tool layer rather than via prompts; the explicit idea of 'MCP as a governance layer' sitting between the agent and downstream systems; input/output filters and RAG grounding to reduce injection and hallucination; and risk-based human oversight plus monitoring, failsafes and incident reporting.",
  "STANDARD 2 — MAS 'Safeguards for Agentic Finance at Runtime' (SAFR), white paper v1.0 dated 3 July 2026, supplemented by MAS FEAT principles and the Veritas methodology as analogues. SAFR is the only source purpose-built for RUNTIME, pre-execution governance of agent actions, and it maps almost one-to-one onto our app's tool choke point. Its core is a four-component checkpoint — Agent Identity, Controls Repository, Disposition Engine, Audit Log — that packages each proposed action into a 'Governance Envelope' and resolves it to one of four dispositions: Deny, Escalate, Auto-Execute, or Observe. SAFR is the architectural spine of our Group A design.",
  "MANDATORY vs RECOMMENDED — THE KEY FINDING (bottom panel). None of these instruments is legally binding today. SAFR explicitly states it 'does not constitute regulatory guidance or supervisory expectations'; IMDA frameworks are voluntary; and the MAS AI Risk Management Guidelines are still only in consultation. Because of that, our control catalogue scores every control on TWO axes: Axis A = binding force, which is 'Recommended (advisory)' for ALL controls; Axis B = the source's own normative strength (core / should / may). So when we later say a control is 'mandatory', we mean it is a framework's own core/'must' language — NOT a legal mandate. Genuine legal obligations would only arise from applicable law plus a finalised MAS AIRG, which is out of scope. Presenter takeaway: we treat these as best-practice engineering priorities, not compliance obligations.",
])

set_notes(s2, [
  "PURPOSE OF THIS SLIDE: explain HOW we structure the defence. The controls from the standards are organised into four control groups (A, B, C, D). Together they form a defence-in-depth control plane; each group attaches to a specific seam in the agentic app. This slide replaces any 'phase' framing — the groups are the real structure.",
  "GROUP A — Action-time authorization. This is the heart of the defence: it decides whether a proposed tool action is allowed to execute, at the tool boundary (mcpCallTool). We build the SAFR checkpoint here — a governance envelope built from trusted state, verified agent identity, a machine-readable mandate of what each agent may do, a deterministic disposition (Deny/Escalate/Auto-Execute/Observe), and exposure + rate limits. It aligns to the MAS SAFR four-component checkpoint and IMDA least-privilege. It protects against the highest-impact threats: fraudulent payouts, hijacked or unauthorized actions, and action abuse (e.g. mass submissions). Model guardrails cannot stop a bad ACTION — only Group A can.",
  "GROUP B — Model input/output guardrails. This screens what enters and leaves the LLM/VLM, at the model hooks. We build prompt-injection detection on chat text AND on the receipt image (treated as untrusted input), PII redaction, and evidence-grounded output checks. It aligns to IMDA GenAI input filters and grounded-output/RAG guidance. It protects against agent hijack via a malicious receipt or injection, and PII leakage. B stops the hijack from landing; A contains the agent if it does. They are complementary — neither replaces the other.",
  "GROUP C — Human oversight and failsafes. This keeps a human in control of high-risk actions, at the humanEscalation node. We build risk-calibrated escalation (the disposition engine decides what needs a human), a substantive escalation contract with a timeout that defaults to deny, a fail-closed floor, and an employee recourse/appeal path. It aligns to SAFR's human-reviewer escalation and IMDA risk-based oversight. It protects against unreviewed high-risk or irreversible actions and rubber-stamp reviews (e.g. 'no reply' silently becoming approval).",
  "GROUP D — Audit, monitoring and incident. This provides a tamper-evident record and the ability to detect and respond, at the audit log. We build an immutable, PII-safe audit log, real-time monitoring wired to actual interventions (halt/terminate/fallback — not just dashboards), and an incident workflow. It aligns to the SAFR Audit Log and IMDA tamper-evident logging/monitoring. It protects against audit tampering, undetected attacks, and forensic gaps.",
  "WHAT WE ARE BUILDING NOW (footer). We start with Group A because it has the highest attack-defense value and, per the gap analysis, 83% of it is absent today. We are building an app-agnostic, in-process control plane at the tool boundary that uses DETERMINISTIC checks against forge-proof trusted state — there is no LLM in the decision loop, which matters because an LLM judge would itself be injectable. The Group A POC is delivered as six thin, independently testable slices: envelope → deny-unknown-tool → identity + mandate → integrity → exposure/rate/evidence → schema + escalate. Groups B, C and D follow incrementally on the same control plane.",
])

set_notes(s3, [
  "PURPOSE OF THIS SLIDE: show where the Expense AI app stands TODAY against the recommended controls, from an evidence-based code audit (every score is backed by a specific file/line reading). This is the baseline we are fixing — it is NOT our build progress.",
  "HEADLINE: strong seams, near-zero enforcement. The app already has the right interception points (the mcpCallTool choke point, model hooks, a humanEscalation node, structured logging), but none of them actually governs. Of 28 assessable controls, ZERO are enforced today.",
  "SCORING KEY (legend) — read this carefully because it explains why 'Present' is 0. Present = a control is built AND deterministically enforced. Partial = a mechanism exists but is prompt-based, incomplete, or non-enforcing (e.g. a system-prompt instruction, a model-asserted check, or a log-only redaction). Absent = nothing found. Scored today: 0 Present, 11 Partial, 17 Absent. Overall that is 61% of controls missing entirely (Absent) and 0% enforced. The 11 'Partial' items are NOT real controls — they are hopeful mechanisms an attacker can bypass.",
  "GROUP A — Action-time authorization: 83% missing (10 of 12 absent, 2 partial, 0 enforced). What exists is only typed MCP schemas, a SELECT-only query guard, and a confidence threshold that is dead config (defined but never read). The biggest gap: mcpCallTool performs NO authorization — no envelope, identity, mandate, disposition, or value/rate limits. This is the single highest-leverage gap.",
  "GROUP B — Model input/output guardrails: 33% missing (2 of 6 absent, 4 partial, 0 enforced). Grounding, PII handling and graceful-failure exist but are prompt-based, model-asserted, or log-only. The biggest gap: no prompt-injection detection at all, and the receipt image is trusted as input — the classic agentic attack vector.",
  "GROUP C — Human oversight: 20% missing (1 of 5 absent, 4 partial, 0 enforced). Reviewers DO hold real authority and escalation persists a status, so this is the strongest group. But there is no timeout, no deadline, no default-deny, and no oversight-effectiveness metrics — an escalation can sit forever.",
  "GROUP D — Audit and monitoring: 80% missing (4 of 5 absent, 1 partial, 0 enforced). An end-to-end trace exists across logs and the LangGraph checkpointer, but the audit_log table is mutable and CASCADE-deletes with the claim, and there is no real-time monitoring wired to intervention and no incident flow.",
  "CONCLUSION (bottom bar): overall 61% of controls are missing entirely and 0% are enforced; Group A is the worst at 10 of 12 absent, so the POC targets Group A first — front-loading the fail-closed deny of unauthorized high-impact actions.",
])

set_notes(s4, [
  "PURPOSE OF THIS SLIDE: explain what we are actually building in the agentic-governance repository — the control-plane design — and how it is engineered to be agnostic to any specific agentic app, while still integrating cleanly with the Expense app.",
  "THE THREE LAYERS (left diagram, top to bottom). Layer 3 — the Host agentic app (e.g. Expense AI) keeps its logic unchanged; its only responsibilities are to expose trusted state through contextvars (which it already sets: authenticated employee id, the VLM-extracted receipt, the session claim id) and to make ONE minimal dependency-injection edit at its composition root so tool calls route through governance. Layer 2 — the Integration adapter is the only app-aware glue: it provides the governedMcpCallTool wrapper (identical async signature to the app's real tool boundary) and a CallContext assembler. Crucially it is INJECTED, not imported — it receives the trusted-state providers and the real tool via an install() call, so governance never imports the app. Layer 1 — the Governance CORE is fully app-agnostic and imports nothing app-specific.",
  "WHAT IS IN THE CORE. A deterministic decision engine that returns exactly one of Deny / Escalate / Auto-Execute / Observe, operating on a Governance Envelope built from trusted state. A set of Ports (abstract interfaces): identity registry, mandate store, policy-decision-point, counters, evidence evaluator, and audit sink. And pluggable Adapters implementing those ports: a pure-Python policy engine (chosen for the POC; OPA is a documented later drop-in), a JSONL audit sink, and in-memory registry and counters. Because the decision logic is deterministic, it is testable, auditable, and cannot be talked around — there is deliberately no LLM in the loop.",
  "HOW IT STAYS APP-AGNOSTIC (right column). First, core/ and ports/ import nothing app-specific — verified in code review (stdlib only). Second, the app injects its providers and real tool via install(), so the core has zero dependency on the app. Third, all policy — allowlists, mandates, thresholds — lives in the governance repo, not the app. Fourth, there is exactly one integration seam today (the tool boundary, mcpCallTool); model hooks are added later for Group B.",
  "ONBOARDING A NEW APP (green box). To govern a different agentic application you only do two things: (1) write a thin adapter that maps that app's ambient state to a CallContext, and (2) add one composition-root DI edit to call the wrapper. The core and the policy engine are unchanged. The Expense app is simply the first adapter.",
  "BOTTOM LINE (dark footer): decisions are deterministic, there is no LLM in the loop, and the system fails closed — if governance is unavailable, high-impact actions are denied rather than allowed through. This is what makes the layer both reusable and safe by default.",
])

out = "docs/slides/agentic-governance-deck.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
